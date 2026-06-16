import sys
from pptx import Presentation

def inspect(file_path):
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            print(f"--- Slide {i+1} ---")
            for j, shape in enumerate(slide.shapes):
                shape_info = f"  Shape {j} ({shape.name}, Type: {shape.shape_type})"
                if shape.has_text_frame:
                    text = shape.text.replace('\n', ' | ').strip()
                    print(f"{shape_info} -> Text: {text}")
                else:
                    print(shape_info)
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    inspect("C:\\Users\\HARDIK\\Downloads\\Presentation_PPt.pptx")
