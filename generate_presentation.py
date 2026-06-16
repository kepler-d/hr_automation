import sys
from pptx import Presentation

def create_presentation():
    input_file = r"C:\Users\HARDIK\Downloads\Presentation_PPt.pptx"
    output_file = r"C:\Users\HARDIK\Downloads\HR_Automation_Final_Presentation.pptx"

    prs = Presentation(input_file)

    # Slide 1: Cover Page
    cover_text = "Project Title: HR Automation Platform\nSubmitted By: Hardik Derashri"
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        if slide.shapes and slide.shapes[0].has_text_frame:
            slide.shapes[0].text = cover_text

    # Slide 2: Table of Contents
    toc_text = (
        "Slide 1: Title & Project Abstract\n"
        "Slide 2: Problem Statement & Engineering Context\n"
        "Slide 3: Pipeline System Architecture\n"
        "Slide 4: Data Engineering Pipeline\n"
        "Slide 5: Challenges & Solutions\n"
        "Slide 6: Containerization & Model Serving\n"
        "Slide 7: Monitoring UI & Client Interface\n"
        "Slide 8: Core Technical Challenges & Debugging\n"
        "Slide 9: Development & Iteration Roadmap"
    )
    if len(prs.slides) > 1:
        slide = prs.slides[1]
        if slide.shapes and slide.shapes[0].has_text_frame:
            slide.shapes[0].text = toc_text

    # Data for Content Slides
    content_slides = [
        # Content 1
        "Title & Project Abstract\n"
        "• Project: TalentFlow AI - HR Automation Platform\n"
        "• Abstract: An intelligent microservices-based application designed to fully automate the resume screening process.\n"
        "• Goal: Leverage Large Language Models (Google Gemini AI) to instantly analyze, score, and rank candidates against job descriptions.\n"
        "• Impact: Reduces manual HR screening time by 90%, eliminating bias and identifying top talent with high accuracy.",
        
        # Content 2
        "Problem Statement & Engineering Context\n"
        "• The Problem: Recruiters spend countless hours manually reading hundreds of resumes per job posting, leading to fatigue, bias, and missed talent.\n"
        "• Context: Traditional ATS (Applicant Tracking Systems) rely on rigid keyword matching, failing to understand semantic context.\n"
        "• Solution: We built a semantic AI screener that reads resumes like a human, understanding contextual skills, experience, and education, providing a strictly calculated score out of 100.",
        
        # Content 3
        "Pipeline System Architecture\n"
        "• Microservices Architecture deployed across multiple specialized cloud providers.\n"
        "• Frontend: React.js and Vite, hosted globally on Vercel for fast CDN delivery.\n"
        "• Backend: Node.js and Express API, hosted on Render for reliable compute.\n"
        "• Database: PostgreSQL, hosted on Neon for serverless relational data storage.\n"
        "• Automations: n8n workflow engine, hosted on Railway for backend event triggers.",
        
        # Content 4
        "Data Engineering Pipeline\n"
        "• Step 1: Bulk File Upload - The React client bundles multiple PDF/TXT files into a single FormData payload.\n"
        "• Step 2: In-Memory Parsing - The Node backend uses 'pdf-parse' to extract text.\n"
        "• Step 3: Fallback OCR - If a PDF is an image, the system automatically falls back to Tesseract.js for Optical Character Recognition.\n"
        "• Step 4: AI Extraction - Text is passed to Gemini 2.5 Flash via a strict prompt to extract Name, Email, Skills, and calculate a score.",
        
        # Content 5
        "Challenges & Solutions\n"
        "• Challenge: API Rate Limits & Timeouts.\n"
        "• Solution: Re-engineered the frontend to send bulk files in one payload, allowing the backend to sequence AI requests and avoid Vercel's strict 10-second timeouts.\n"
        "• Challenge: Heavy Local AI Models.\n"
        "• Solution: Initially built with Ollama/Llama3 locally, which consumed massive RAM. Migrated to the cloud-native Google Gemini API for exponential speed improvements.",
        
        # Content 6
        "Containerization & Model Serving\n"
        "• Docker Configuration: The entire backend is containerized using Docker and Docker Compose.\n"
        "• Multi-stage Builds: Optimized Dockerfile to keep the Node.js production image lightweight and fast.\n"
        "• Model Serving: Transitioned from local containerized LLMs to highly-available REST API endpoints (Gemini SDK) to ensure production-grade scalability.",
        
        # Content 7
        "Monitoring UI & Client Interface\n"
        "• Tech Stack: React, Tailwind CSS, Google Material Symbols.\n"
        "• Features: Dark mode design, interactive dashboard, drag-and-drop file screening.\n"
        "• State Management: Real-time filtering, instant candidate status updates (Pending, Shortlisted, Rejected).\n"
        "• Aesthetics: High-quality glassmorphism UI, interactive hover states, and smooth score progress bars.",
        
        # Content 8
        "Core Technical Challenges & Debugging\n"
        "• PostgreSQL Inserts: SQLite automatically returns IDs on insert, but PostgreSQL requires explicit '.returning()' clauses. Debugged and resolved this Knex.js dialect difference.\n"
        "• Syntax Errors: Resolved backend crash loops caused by duplicated variable declarations during the massive AI migration.\n"
        "• Frontend Sequencing: Debugged an issue where only one resume appeared in bulk uploads by consolidating multiple async fetch requests into a unified batch.",
        
        # Content 9
        "Development & Iteration Roadmap\n"
        "• Future Additions:\n"
        "• WebSocket Integration: To stream real-time progress updates to the UI during bulk AI processing.\n"
        "• Webhook Automations: Triggering n8n to automatically email candidates their interview links once their status is changed to 'Shortlisted'.\n"
        "• Advanced Schemas: Implementing strict JSON Schema validation in the Gemini API to guarantee 100% reliable data structures."
    ]

    # Use Slide 3's layout for all content slides
    layout = prs.slides[2].slide_layout

    # Update existing slides 3, 4, 5
    for i in range(2, 5):
        if i < len(prs.slides):
            slide = prs.slides[i]
            if slide.shapes and slide.shapes[0].has_text_frame:
                slide.shapes[0].text = content_slides[i - 2]
    
    # Append the remaining slides
    for i in range(3, len(content_slides)):
        slide = prs.slides.add_slide(layout)
        if slide.shapes and slide.shapes[0].has_text_frame:
            slide.shapes[0].text = content_slides[i]

    prs.save(output_file)
    print(f"Presentation saved to: {output_file}")

if __name__ == "__main__":
    create_presentation()
